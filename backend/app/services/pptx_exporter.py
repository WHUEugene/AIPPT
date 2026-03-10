from __future__ import annotations

from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Optional

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Inches

from ..schemas.project import ProjectState
from ..utils.logger import get_logger


class PPTXExporter:
    DEFAULT_SLIDE_HEIGHT_INCHES = 7.5
    EXPORT_DPI = 180
    JPEG_QUALITY = 90

    def __init__(self, output_dir: Path, image_dir: Path):
        self.output_dir = Path(output_dir)
        self.image_dir = Path(image_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.logger = get_logger()

    def _get_image_dimensions(self, image_path: Path) -> tuple[int, int]:
        with Image.open(image_path) as img:
            return img.size

    def _parse_aspect_ratio(self, aspect_ratio: str | None) -> tuple[int, int] | None:
        if not aspect_ratio or ":" not in aspect_ratio:
            return None
        try:
            width_ratio, height_ratio = (
                int(part.strip()) for part in aspect_ratio.split(":", 1)
            )
        except ValueError:
            return None
        if width_ratio <= 0 or height_ratio <= 0:
            return None
        return width_ratio, height_ratio

    def _resolve_slide_ratio(self, project: ProjectState) -> tuple[int, int]:
        for slide in project.slides:
            image_path = self._resolve_image_path(slide.image_url)
            if image_path and image_path.exists():
                return self._get_image_dimensions(image_path)

        ratio = self._parse_aspect_ratio(project.aspect_ratio)
        if ratio is not None:
            return ratio

        return (16, 9)

    def _slide_dimensions(self, project: ProjectState) -> tuple[int, int]:
        width_ratio, height_ratio = self._resolve_slide_ratio(project)
        slide_height_inches = self.DEFAULT_SLIDE_HEIGHT_INCHES
        slide_width_inches = slide_height_inches * width_ratio / height_ratio
        return int(Inches(slide_width_inches)), int(Inches(slide_height_inches))

    def _target_pixel_size(self, slide_width: int, slide_height: int) -> tuple[int, int]:
        width_px = max(1, int(round((slide_width / 914400) * self.EXPORT_DPI)))
        height_px = max(1, int(round((slide_height / 914400) * self.EXPORT_DPI)))
        return width_px, height_px

    def _prepare_image_blob(
        self,
        image_path: Path,
        slide_width: int,
        slide_height: int,
    ) -> tuple[BytesIO, dict[str, str | int]]:
        with Image.open(image_path) as img:
            img = ImageOps.exif_transpose(img)
            original_width, original_height = img.size
            target_width_px, target_height_px = self._target_pixel_size(slide_width, slide_height)
            resampling = Image.Resampling.LANCZOS if hasattr(Image, "Resampling") else Image.LANCZOS
            if img.mode != "RGB":
                flattened = Image.new("RGB", img.size, "white")
                alpha = img.getchannel("A") if "A" in img.getbands() else None
                flattened.paste(img, mask=alpha)
                img = flattened

            if img.size != (target_width_px, target_height_px):
                img = img.resize((target_width_px, target_height_px), resampling)

            output = BytesIO()
            img.save(
                output,
                format="JPEG",
                quality=self.JPEG_QUALITY,
                optimize=True,
                subsampling=0,
            )
            output.seek(0)

        return output, {
            "original_size": f"{original_width}x{original_height}",
            "export_size": f"{target_width_px}x{target_height_px}",
            "export_bytes": output.getbuffer().nbytes,
        }

    def _add_picture_with_protection(
        self,
        slide,
        image_path: Path,
        slide_width: int,
        slide_height: int,
        session_id: str,
    ) -> bool:
        try:
            image_blob, metadata = self._prepare_image_blob(image_path, slide_width, slide_height)
            slide.shapes.add_picture(
                image_blob,
                0,
                0,
                width=int(slide_width),
                height=int(slide_height),
            )
            self.logger.log_pipeline_step(
                session_id=session_id,
                step="image_added",
                details={
                    "image_file": image_path.name,
                    "original_size": metadata["original_size"],
                    "export_size": metadata["export_size"],
                    "export_bytes": metadata["export_bytes"],
                    "slide_size_emu": f"{int(slide_width)}x{int(slide_height)}",
                    "stage": "图片按原始比例铺满导出页面后已添加",
                },
            )
            return True
        except Exception as exc:
            self.logger.logger.error(f"Failed to add image to slide: {exc}")
            return False

    def build(self, project: ProjectState, session_id: Optional[str] = None) -> tuple[str, BytesIO]:
        prs = Presentation()

        if not session_id:
            session_id = self.logger.start_session(
                "pptx_export",
                project_title=project.title,
                slides_count=len(project.slides),
            )

        slide_width, slide_height = self._slide_dimensions(project)
        prs.slide_width = slide_width
        prs.slide_height = slide_height

        self.logger.log_pipeline_step(
            session_id=session_id,
            step="slide_size_adjusted",
            details={
                "aspect_ratio": project.aspect_ratio or "auto",
                "slide_width_emu": int(slide_width),
                "slide_height_emu": int(slide_height),
                "slide_width_inches": slide_width / 914400,
                "slide_height_inches": slide_height / 914400,
            },
        )

        images_processed = 0
        images_failed = 0

        try:
            for i, slide_data in enumerate(project.slides):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                image_path = self._resolve_image_path(slide_data.image_url)

                self.logger.log_pipeline_step(
                    session_id=session_id,
                    step="slide_image_check",
                    details={
                        "slide_index": i + 1,
                        "slide_title": slide_data.title,
                        "image_url": slide_data.image_url,
                        "resolved_path": str(image_path) if image_path else "None",
                        "path_exists": image_path.exists() if image_path else False,
                        "stage": f"第{i + 1}张幻灯片图片路径检查",
                    },
                )

                if image_path and image_path.exists():
                    success = self._add_picture_with_protection(
                        slide,
                        image_path,
                        slide_width,
                        slide_height,
                        session_id,
                    )
                    self.logger.log_pipeline_step(
                        session_id=session_id,
                        step="slide_image_add_result",
                        details={
                            "slide_index": i + 1,
                            "slide_title": slide_data.title,
                            "success": success,
                            "image_path": str(image_path),
                            "stage": f"第{i + 1}张幻灯片图片添加结果: {'成功' if success else '失败'}",
                        },
                    )
                    if success:
                        images_processed += 1
                    else:
                        images_failed += 1
                else:
                    self.logger.log_pipeline_step(
                        session_id=session_id,
                        step="slide_no_image",
                        details={
                            "slide_index": i + 1,
                            "slide_title": slide_data.title,
                            "stage": f"第{i + 1}张幻灯片无图片",
                        },
                    )

            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            filename = self._filename(project)

            self.logger.log_response(
                session_id=session_id,
                stage="pptx_export_complete",
                data={
                    "filename": filename,
                    "total_slides": len(project.slides),
                    "images_processed": images_processed,
                    "images_failed": images_failed,
                    "success_rate": images_processed / len(project.slides) if project.slides else 0,
                    "pptx_bytes": buffer.getbuffer().nbytes,
                },
                success=True,
            )

            return filename, buffer

        except Exception as exc:
            self.logger.log_response(
                session_id=session_id,
                stage="pptx_export_error",
                data={
                    "error": str(exc),
                    "error_type": type(exc).__name__,
                    "images_processed": images_processed,
                    "images_failed": images_failed,
                },
                success=False,
            )
            raise

    def _resolve_image_path(self, image_url: str | None) -> Path | None:
        if not image_url:
            return None

        if "://" in image_url:
            if "assets/" in image_url:
                file_name = image_url.split("assets/", 1)[1]
                return self.image_dir / file_name
            return None

        sanitized = image_url.lstrip("/")
        if sanitized.startswith("assets/"):
            file_name = sanitized.split("assets/", 1)[1]
            return self.image_dir / file_name
        if (self.image_dir / sanitized).exists():
            return self.image_dir / sanitized
        return Path(sanitized)

    def _filename(self, project: ProjectState) -> str:
        timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        base = project.title or "AI_PPT_Flow"
        safe = "".join(char for char in base if char.isalnum() or char in ("_", "-"))[:40] or "AI_PPT_Flow"
        return f"{safe}_{timestamp}.pptx"


__all__ = ["PPTXExporter"]
